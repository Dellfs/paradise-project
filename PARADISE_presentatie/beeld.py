# -*- coding: utf-8 -*-
"""Getekende beelden: de drukmat en het doorlopende trajectcanvas.

Beide zijn opgebouwd uit echte PowerPoint-vormen, dus ze blijven bewerkbaar
en — belangrijker — ze kunnen morphen. Elke vorm die over dia's heen moet
morphen krijgt een naam die met !! begint; dat is de enige manier om
PowerPoint te dwingen twee vormen aan elkaar te koppelen in plaats van te
gokken (Microsoft, "Morph transition: tips and tricks").
"""
import math

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from inhoud import K, FONT, FONT_M

B = 1920
SCH = 13.333 / B


def px(v):
    return Inches(v * SCH)


def rgb(n):
    return RGBColor.from_string(K[n] if n in K else n)


# --------------------------------------------------------------- voetprofiel
# Mediale en laterale rand van een rechtervoet, van hiel (0) naar teen (1).
# Grof maar herkenbaar: de inham op 0,34-0,45 is de mediale voetboog.
PROFIEL = [
    (0.00, 0.40, 0.60), (0.03, 0.33, 0.67), (0.08, 0.29, 0.71),
    (0.15, 0.28, 0.72), (0.24, 0.29, 0.73), (0.32, 0.33, 0.74),
    (0.40, 0.37, 0.75), (0.48, 0.36, 0.77), (0.55, 0.31, 0.79),
    (0.62, 0.24, 0.81), (0.69, 0.19, 0.82), (0.75, 0.17, 0.81),
    (0.81, 0.17, 0.78), (0.86, 0.18, 0.74), (0.90, 0.19, 0.70),
    (0.94, 0.20, 0.64), (0.97, 0.21, 0.58), (1.00, 0.23, 0.50),
]


def _rand(t):
    """Mediale en laterale rand op hoogte t (0 = hiel, 1 = teen)."""
    if t <= 0:
        return PROFIEL[0][1], PROFIEL[0][2]
    if t >= 1:
        return PROFIEL[-1][1], PROFIEL[-1][2]
    for i in range(len(PROFIEL) - 1):
        t0, l0, r0 = PROFIEL[i]
        t1, l1, r1 = PROFIEL[i + 1]
        if t0 <= t <= t1:
            f = (t - t0) / (t1 - t0)
            return l0 + f * (l1 - l0), r0 + f * (r1 - r0)
    return PROFIEL[-1][1], PROFIEL[-1][2]


# ------------------------------------------------------------------ kleurband
# Dit is de kleurband uit het PARADISE-logo zelf: de voetafdruk in het merk is
# een drukmat, van diep blauw via cyaan en geel naar oranje en rood. De schaal
# van de dia's en het logo zijn dus letterlijk dezelfde.
BAND = [(0.00, (0x05, 0x1D, 0x33)), (0.16, (0x0D, 0x3A, 0x63)),
        (0.32, (0x1D, 0x8D, 0xB0)), (0.48, (0x52, 0xBD, 0xEC)),
        (0.62, (0xFF, 0xD2, 0x4A)), (0.80, (0xFF, 0x7A, 0x00)),
        (1.00, (0xE8, 0x33, 0x1E))]


def kleur(v):
    v = max(0.0, min(1.0, v))
    for i in range(len(BAND) - 1):
        a, ca = BAND[i]
        b, cb = BAND[i + 1]
        if a <= v <= b:
            f = 0 if b == a else (v - a) / (b - a)
            return '%02X%02X%02X' % tuple(
                int(round(ca[j] + f * (cb[j] - ca[j]))) for j in range(3))
    return '%02X%02X%02X' % BAND[-1][1]


def veld(hotspots, x, y):
    """Som van gaussische drukhaarden op genormaliseerde voetcoördinaten."""
    v = 0.0
    for hx, hy, amp, sig in hotspots:
        d2 = (x - hx) ** 2 + (y - hy) ** 2
        v += amp * math.exp(-d2 / (2 * sig * sig))
    return v


# De hele zool draagt iets; de haarden liggen daar bovenop. De piek ligt op
# metatarsaal 2-3, de klassieke plek van een plantair ulcus.
BASIS = (0.48, 0.46, 0.26, 0.40)
VOOR = [BASIS, (0.50, 0.10, 0.60, 0.105), (0.44, 0.66, 1.00, 0.062),
        (0.30, 0.68, 0.55, 0.058), (0.70, 0.63, 0.42, 0.062),
        (0.28, 0.91, 0.46, 0.055), (0.60, 0.40, 0.12, 0.110)]
# Na optimalisatie: de piek daalt, de voetboog gaat mee dragen.
NA = [BASIS, (0.50, 0.12, 0.54, 0.115), (0.44, 0.66, 0.50, 0.080),
      (0.30, 0.68, 0.44, 0.070), (0.70, 0.63, 0.40, 0.075),
      (0.28, 0.91, 0.40, 0.062), (0.58, 0.42, 0.42, 0.130)]


def drukmat(s, x, y, w, h, hotspots, kol=18, rij=38, sleutel='a',
            piek=312.0, schaal_max=400.0):
    """Sensormatrix in de vorm van een voet — wat een pedar-zool oplevert.

    De waarden worden zo herschaald dat de heetste cel exact op de opgegeven
    piekdruk uitkomt; anders klopt het beeld niet met het cijfer ernaast.

    Elke cel heet !!<sleutel><rij>_<kolom>, zodat morph de cel van de ene dia
    aan dezelfde cel op de volgende koppelt en alleen de kleur laat bewegen.
    """
    cw, ch = w / kol, h / rij
    d = min(cw, ch) * 0.86
    cellen = []
    for r in range(rij):
        t = (r + 0.5) / rij
        li, la = _rand(t)
        for c in range(kol):
            fx = (c + 0.5) / kol
            if li <= fx <= la:
                cellen.append((r, c, veld(hotspots, fx, t)))
    top = max(v for _, _, v in cellen) or 1.0
    f = (piek / schaal_max) / top
    for r, c, v in cellen:
        cx = x + c * cw + (cw - d) / 2
        cy = y + h - (r + 1) * ch + (ch - d) / 2
        o = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               px(cx), px(cy), px(d), px(d))
        o.fill.solid()
        o.fill.fore_color.rgb = rgb(kleur(v * f))
        o.line.fill.background()
        o.shadow.inherit = False
        o.adjustments[0] = 0.3
        o.name = '!!%s%02d_%02d' % (sleutel, r, c)
    return len(cellen)


def schaalbalk(s, x, y, w, h, txt, maxwaarde=400):
    """Legende onder de drukmat: van koud naar heet, met kPa-schaal."""
    stappen = 60
    for i in range(stappen):
        o = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               px(x + i * w / stappen), px(y),
                               px(w / stappen + 1), px(h))
        o.fill.solid()
        o.fill.fore_color.rgb = rgb(kleur(i / (stappen - 1.0)))
        o.line.fill.background()
        o.shadow.inherit = False
    for f in (0, 0.5, 1.0):
        txt(s, x + f * w - 60, y + h + 10, 120, 30,
            '%d' % round(f * maxwaarde), gr=11, kl='gedempt', font=FONT_M,
            uit=PP_ALIGN.CENTER)
    txt(s, x + w + 24, y - 2, 200, 30, 'kPa', gr=11, kl='gedempt', font=FONT_M)


# ------------------------------------------------------- doorlopend trajectcanvas
# Wereldcoördinaten vallen bij schaal 1 samen met de dia. Elke zoomdia tekent
# hetzelfde canvas met een andere schaal en een ander middelpunt; morph maakt
# daar een camerabeweging van.
W_LINKS, W_BREED = 100, 1720
W_AS, W_LABEL, W_TAAK = 454, 330, 506

# De stip zegt wat voor moment het is, niet of het extra is.
KLEUR_SOORT = {'meting': 'licht', 'zorg': 'mid', 'start': 'oranje', 'eind': 'oranje'}


def legende(s, x, y, soorten, txt):
    """Drie stippen met een label: wat betekent welke kleur op de tijdlijn."""
    for i, (soort, wat) in enumerate(soorten):
        cx = x + i * 400
        o = s.shapes.add_shape(MSO_SHAPE.OVAL, px(cx), px(y), px(18), px(18))
        o.fill.solid(); o.fill.fore_color.rgb = rgb(KLEUR_SOORT[soort])
        o.line.fill.background(); o.shadow.inherit = False
        txt(s, cx + 30, y - 6, 360, 34, wat, gr=13, kl='gedempt', font=FONT_M)


def station_x(j, n):
    stap = W_BREED / n
    return W_LINKS + stap * j + stap / 2


def canvas(s, punten, txt, liniaal, schaal=1.0, cx=960.0, cy=540.0,
           ox=960.0, oy=540.0):
    """Tekent het trajectcanvas onder een camera-instelling.

    schaal 1 met cx/cy = ox/oy geeft het overzicht; een hogere schaal met een
    ander middelpunt zoomt in op één fase. Alles behoudt zijn !!-naam, dus
    PowerPoint beweegt de camera in plaats van de dia te verversen.
    """
    def X(wx):
        return (wx - cx) * schaal + ox

    def Y(wy):
        return (wy - cy) * schaal + oy

    n = len(punten)
    stap = W_BREED / n
    lijn = liniaal(s, X(W_LINKS), Y(W_AS), W_BREED * schaal, 'rand',
                   max(2.0, 3 * schaal))
    lijn.name = '!!as'

    for j, p in enumerate(punten):
        wanneer, wat, soort = p[0], p[1], p[2]
        kl = KLEUR_SOORT.get(soort, 'mid')
        wx = station_x(j, n)
        dm = (30 if soort in ('start', 'eind') else 24) * schaal
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, px(X(wx) - dm / 2),
                               px(Y(W_AS + 2) - dm / 2), px(dm), px(dm))
        c.fill.solid(); c.fill.fore_color.rgb = rgb(kl)
        c.line.color.rgb = rgb('tegel'); c.line.width = Pt(3 * schaal)
        c.shadow.inherit = False
        c.name = '!!stip%d' % j

        t1 = txt(s, X(wx - stap / 2), Y(W_LABEL), stap * schaal, 44 * schaal,
                 wanneer, gr=(15 if len(wanneer) > 4 else 17) * schaal, kl='ink',
                 vet=True, uit=PP_ALIGN.CENTER, font=FONT_M)
        t1.name = '!!maand%d' % j
        if wat:
            t3 = txt(s, X(wx - stap / 2 + 10), Y(W_TAAK), (stap - 20) * schaal,
                     90 * schaal, wat, gr=13.5 * schaal, kl='gedempt',
                     uit=PP_ALIGN.CENTER, ra=1.3)
            t3.name = '!!taak%d' % j
