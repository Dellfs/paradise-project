# -*- coding: utf-8 -*-
"""Bouwt de PARADISE-opleidingssessie: dark mode, bento-grid, morph-overgangen.

Alle tekst staat in echte tekstvakken, dus alles blijft bewerkbaar in
PowerPoint. Sprekersnotities bevatten per dia de presentatietip en het
interactiemoment. Inhoud staat in inhoud.py.
"""
import os, sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image

from inhoud import SLIDES, K, FONT, FONT_L, FONT_M, kol, PUNTEN, SOORT, DECKS

# Welk deck bouwen we? `python maak_pptx.py board` of `... alle`.
KEUZE = (sys.argv[1] if len(sys.argv) > 1 else 'opleiding').lower()
import beeld

HIER = os.path.dirname(os.path.abspath(__file__))
BEELD = os.path.join(HIER, 'beeld')

B, H = 1920, 1080
SCH = 13.333 / B


def px(v):
    return Inches(v * SCH)


def rgb(n):
    return RGBColor.from_string(K[n])


def tegel(s, x, y, w, h, vul='tegel', rand='rand', alpha=None, rond=True, naam=None):
    """Bento-tegel; alpha geeft het glaseffect."""
    vorm = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rond else MSO_SHAPE.RECTANGLE,
        px(x), px(y), px(w), px(h))
    if naam:
        vorm.name = naam
    vorm.fill.solid()
    vorm.fill.fore_color.rgb = rgb(vul)
    if alpha is not None:
        sf = vorm.fill.fore_color._xFill.find(qn('a:srgbClr'))
        el = etree.SubElement(sf, qn('a:alpha'))
        el.set('val', str(int(alpha * 1000)))
    if rand:
        vorm.line.color.rgb = rgb(rand)
        vorm.line.width = Pt(0.75)
    else:
        vorm.line.fill.background()
    if rond:
        vorm.adjustments[0] = 0.045
    vorm.shadow.inherit = False
    return vorm


def txt(s, x, y, w, h, regels, gr=16, kl='ink', vet=False, font=FONT,
        uit=PP_ALIGN.LEFT, ra=1.25, sp=0, anker=MSO_ANCHOR.TOP, caps=False, naam=None,
        omslag=True, vloei=None):
    tb = s.shapes.add_textbox(px(x), px(y), px(w), px(h))
    if naam:
        tb.name = naam
    tf = tb.text_frame
    tf.word_wrap = omslag
    tf.vertical_anchor = anker
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(regels, str):
        regels = regels.split('\n')
    for i, r in enumerate(regels):
        ex = {}
        if isinstance(r, tuple):
            r, ex = r
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ex.get('uit', uit)
        p.line_spacing = ex.get('ra', ra)
        if i:
            p.space_before = Pt(ex.get('voor', 5))
        run = p.add_run()
        run.text = r.upper() if ex.get('caps', caps) else r
        f = run.font
        f.name = ex.get('font', font)
        f.size = Pt(ex.get('gr', gr))
        f.bold = ex.get('vet', vet)
        f.color.rgb = rgb(ex.get('kl', kl))
        s2 = ex.get('sp', sp)
        if s2:
            run.font._rPr.set('spc', str(int(s2 * 100)))
        vk = ex.get('vk', vloei)
        if vk:
            # kleurverloop dwars door de letters, in de gradiënt van het merk
            rPr = run.font._rPr
            oud = rPr.find(qn('a:solidFill'))
            if oud is not None:
                rPr.remove(oud)
            rPr.insert(0, etree.fromstring(
                '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/'
                '2006/main"><a:gsLst>'
                '<a:gs pos="0"><a:srgbClr val="%s"/></a:gs>'
                '<a:gs pos="100000"><a:srgbClr val="%s"/></a:gs>'
                '</a:gsLst><a:lin ang="2700000" scaled="0"/></a:gradFill>'
                % (K[vk[0]], K[vk[1]])))
    return tb


def punt(s, x, y, d, kl):
    """Losse stip voor de puntenmatrix."""
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y), px(d), px(d))
    o.fill.solid(); o.fill.fore_color.rgb = rgb(kl)
    o.line.fill.background(); o.shadow.inherit = False
    return o


def matrix(s, x, y, n, per_rij, d, spatie, kleuren):
    """Puntenmatrix: n stippen, kleur per index via de lijst kleuren."""
    for i in range(n):
        r, k = divmod(i, per_rij)
        punt(s, x + k * (d + spatie), y + r * (d + spatie), d, kleuren[i])


def staaf(s, x, y, w, h, deel, kl, achter='tegel2'):
    """Horizontale staaf met gevuld aandeel (deel = 0..1)."""
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h))
    b.fill.solid(); b.fill.fore_color.rgb = rgb(achter)
    b.line.fill.background(); b.shadow.inherit = False
    b.adjustments[0] = 0.5
    if deel > 0:
        f = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w * deel), px(h))
        f.fill.solid(); f.fill.fore_color.rgb = rgb(kl)
        f.line.fill.background(); f.shadow.inherit = False
        f.adjustments[0] = 0.5
    return b


def merkstreep(s, x, y, h, kl='ink'):
    """Verticale streefwaardestreep over een staaf."""
    m = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(3), px(h))
    m.fill.solid(); m.fill.fore_color.rgb = rgb(kl)
    m.line.fill.background(); m.shadow.inherit = False
    return m


def gloed(s, x, y, d, kl, alpha=7):
    """Grote, bijna doorzichtige cirkel: geeft diepte achter een heldia."""
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y), px(d), px(d))
    o.fill.solid(); o.fill.fore_color.rgb = rgb(kl)
    sf = o.fill.fore_color._xFill.find(qn('a:srgbClr'))
    etree.SubElement(sf, qn('a:alpha')).set('val', str(int(alpha * 1000)))
    o.line.fill.background(); o.shadow.inherit = False
    return o


def hoogte(tekst, breedte, gr, ra=1.3):
    """Schat hoe hoog een stuk tekst wordt, in dia-eenheden.

    Eén eenheid is een halve punt en een letter is in Segoe UI gemiddeld een
    halve em breed. Ruw, maar genoeg om te voorkomen dat blokken over elkaar
    schuiven zodra de tekst langer wordt.
    """
    per_regel = max(1, int((breedte / 2.0) / (gr * 0.5)))
    regels = max(1, -(-len(tekst) // per_regel))
    return regels * gr * ra * 2


def bol(s, x, y, d, tekst, kl='licht', klein=False, vul='tegel2', rand=None):
    """Genummerde bol voor stappen en handelingen.

    Het tekstvak valt samen met de cirkel en is verticaal gecentreerd; met een
    losse verschuiving zakt het cijfer altijd net te laag, omdat de regelafstand
    boven de letter meetelt.
    """
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y), px(d), px(d))
    o.fill.solid(); o.fill.fore_color.rgb = rgb(vul)
    o.line.color.rgb = rgb(rand or kl); o.line.width = Pt(1)
    o.shadow.inherit = False
    txt(s, x, y, d, d, tekst, gr=12.5 if klein else 15, kl=kl, vet=True,
        uit=PP_ALIGN.CENTER, font=FONT_M, ra=1.0, anker=MSO_ANCHOR.MIDDLE)
    return o


def liniaal(s, x, y, w, kl='rand', dik=1.5):
    """Haarlijn als scheiding — vervangt een kader waar een kader te veel is."""
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(dik))
    r.fill.solid(); r.fill.fore_color.rgb = rgb(kl)
    r.line.fill.background(); r.shadow.inherit = False
    return r


def foto(s, bestand, x, y, w, h, naam=None, alpha=None, vullend=False):
    """Plaatst een foto in het vak x,y,w,h.

    vullend=False past hem passend in en centreert; vullend=True laat hem het
    vak volledig vullen en over de randen lopen, zoals een aflopend beeld.
    alpha (0-100) maakt hem doorschijnend, voor beeld dat achtergrond is.
    """
    pad = os.path.join(BEELD, bestand)
    with Image.open(pad) as im:
        bw, bh = im.size
    f = max(w / bw, h / bh) if vullend else min(w / bw, h / bh)
    fw, fh = bw * f, bh * f
    p = s.shapes.add_picture(pad, px(x + (w - fw) / 2), px(y + (h - fh) / 2),
                             px(fw), px(fh))
    if naam:
        p.name = naam
    if alpha is not None:
        blip = p._element.blipFill.find(qn('a:blip'))
        etree.SubElement(blip, qn('a:alphaModFix')).set('amt', str(int(alpha * 1000)))
    return p


# De instellingslogo's zijn merkbestanden die ik niet mag namaken. Zet ze als
# kuleuven.png en vub.png in beeld\ en ze verschijnen automatisch; zolang ze er
# niet zijn, staat de naam er getypt.
PARTNERS = [('kuleuven.png', 'KU Leuven', 'Campus Brugge'),
            ('vub.png', 'VUB', 'Vrije Universiteit Brussel'),
            (None, 'FWO', 'TBM T000226N')]


def partners(s, y=904):
    """Balk met de instellingen, onderaan de openings- en slotdia.

    De logo's staan op wit met donkerblauwe letters, dus ze krijgen een wit
    vlak. Dat is ook hoe beide instellingen hun merk voorschrijven op een
    donkere ondergrond.
    """
    liniaal(s, 100, y - 26, 1720)
    x, hoog = 100, 62
    for bestand, naam, sub in PARTNERS:
        pad = bestand and os.path.join(BEELD, bestand)
        if pad and os.path.exists(pad):
            # het witte vlak past zich aan het logo aan in plaats van omgekeerd
            with Image.open(pad) as im:
                breed = hoog * im.size[0] / float(im.size[1])
            tegel(s, x, y, breed + 48, hoog + 34, 'wit', None)
            foto(s, bestand, x + 24, y + 17, breed, hoog)
            x += breed + 48 + 40
        else:
            txt(s, x, y + 8, 340, 40, naam, gr=21, kl='ink', vet=True)
            txt(s, x + 2, y + 48, 420, 34, sub, gr=12, kl='gedempt', font=FONT_M)
            x += 400


def merk(s, t):
    """Het PARADISE-merk: groot op de openings-, sectie- en slotdia, klein
    rechtsboven op de inhoudsdia's."""
    if t in ('hero_titel', 'sectie', 'knal', 'titel_foto', 'titel_duo'):
        pass  # deze dia's plaatsen het merk zelf, of dragen er geen
    elif t == 'nooit':
        pass  # sectiedia draagt het merk al groot; een knaldia draagt niets
    elif t == 'slot':
        foto(s, 'merk.png', 1380, 320, 330, 350, naam='merk')
    elif t == 'pauze':
        foto(s, 'merk.png', 890, 118, 140, 160, naam='merk')
    else:
        foto(s, 'merk.png', 1742, 58, 76, 76, naam='merk')


def kicker(s, t, kl='licht', y=90):
    txt(s, 100, y, 1500, 34, t, gr=12.5, kl=kl, vet=True, font=FONT_M, sp=2.4, caps=True)


def kop(s, t, y=132, gr=40, kl='ink', w=1560, naam=None):
    txt(s, 100, y, w, 200, t, gr=gr, kl=kl, vet=True, ra=1.1, naam=naam)


def notitie(s, d):
    tf = s.notes_slide.notes_text_frame
    tf.text = ('PRESENTATIETIP\n' + d.get('tip', '') +
               '\n\nINTERACTIE\n' + d.get('interactie', ''))


def morph(s, soort='morph'):
    """PowerPoint-morphovergang, met fade als terugval op oudere versies."""
    if soort == 'fade':
        xml = ('<p:transition xmlns:p="http://schemas.openxmlformats.org/'
               'presentationml/2006/main" xmlns:p14="http://schemas.microsoft.com/'
               'office/powerpoint/2010/main" spd="slow" p14:dur="900"><p:fade/></p:transition>')
    else:
        xml = (
          '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/'
          'markup-compatibility/2006" xmlns:p="http://schemas.openxmlformats.org/'
          'presentationml/2006/main">'
          '<mc:Choice xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main" '
          'Requires="p159">'
          '<p:transition xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
          'spd="slow" p14:dur="1200"><p159:morph option="byObject"/></p:transition></mc:Choice>'
          '<mc:Fallback>'
          '<p:transition xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
          'spd="slow" p14:dur="900"><p:fade/></p:transition></mc:Fallback>'
          '</mc:AlternateContent>')
    s._element.append(etree.fromstring(xml))


def uitklappen(slides, aan=True):
    """Zet elke bezoekdia om in een titelkaart plus één dia per handeling.

    Een keynote loopt op tempo: één gedachte per dia, groot gezet, en doorklikken.
    Alleen het opleidingsdeck heeft die uitwerking nodig; de kortere decks tonen
    het bezoek als één compacte dia. Het veld `kaart` zegt welke van de twee.
    """
    uit = []
    for d in slides:
        if d['t'] == 'bezoek':
            d = dict(d, kaart=aan)
        uit.append(d)
        if d['t'] != 'bezoek' or not aan:
            continue
        n = len(d['handelingen'])
        for j, h in enumerate(d['handelingen']):
            # een handeling mag ook (tekst, fotobijschrift) zijn
            h, bij = h if isinstance(h, tuple) else (h, None)
            uit.append(dict(t='stap', morph='morph', visite=d['kicker'],
                            nr=j + 1, totaal=n, tekst=h, bijschrift=bij,
                            documenten=d['documenten'] if j == n - 1 else [],
                            letop=d['letop'] if j == n - 1 else None,
                            tip=d.get('tip', '') if j == 0 else '',
                            interactie=d.get('interactie', '') if j == 0 else ''))
    return uit


if KEUZE not in DECKS:
    raise SystemExit('Onbekend deck: %s. Kies uit: %s'
                     % (KEUZE, ', '.join(sorted(DECKS))))
DECK = DECKS[KEUZE]
DOEL = os.path.join(HIER, DECK['bestand'])

# alleen de dia's die voor dit deck bedoeld zijn, en de ondertitel op maat
SLIDES = [dict(d) for d in SLIDES if DECK['letter'] in d['voor']]
for d in SLIDES:
    if d['t'].startswith('titel'):
        d['onder'] = DECK['onder']

SLIDES = uitklappen(SLIDES, aan=DECK.get('uitklappen', False))

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
leeg = prs.slide_layouts[6]
TOT = len(SLIDES)


def verloop(vorm, van, naar, hoek=5400000):
    """Lineair kleurverloop op een vorm; hoek in 1/60000 graad."""
    sp = vorm.fill._xPr
    for k in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:blipFill', 'a:pattFill'):
        oud = sp.find(qn(k))
        if oud is not None:
            sp.remove(oud)
    xml = (
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'rotWithShape="1"><a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="%s"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="%s"/></a:gs>'
        '</a:gsLst><a:lin ang="%d" scaled="0"/></a:gradFill>' % (K[van], K[naar], hoek))
    el = etree.fromstring(xml)
    # De volgorde binnen spPr ligt vast: de vulling hoort ná de geometrie.
    geom = sp.find(qn('a:prstGeom'))
    if geom is None:
        geom = sp.find(qn('a:custGeom'))
    if geom is None:
        sp.insert(0, el)
    else:
        geom.addnext(el)


def nieuw():
    """Elke dia krijgt een verlopende grond in plaats van een vlakke kleur:
    diep navy linksboven, iets lichter rechtsonder. Subtiel, maar het haalt de
    dia's uit het vlakke."""
    s = prs.slides.add_slide(leeg)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, px(B), px(H))
    bg.line.fill.background(); bg.shadow.inherit = False
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb('bg')
    verloop(bg, 'grond2', 'bg', 5400000)
    return s


def paginering(s, i):
    txt(s, 1660, 1006, 160, 28, '%02d / %02d' % (i, TOT), gr=10.5, kl='gedempt',
        font=FONT_M, uit=PP_ALIGN.RIGHT)


for i, d in enumerate(SLIDES, 1):
    t, s = d['t'], nieuw()

    if t == 'titel_foto':
        # A — aflopende foto, type rustig onderaan. Het beeld draagt alles.
        foto(s, d.get('foto', 'plaatshouder.png'), 0, 0, 1920, 1080,
             vullend=True, naam='fotokader')
        scrim = tegel(s, 0, 0, 1920, 1080, 'bg', None, rond=False)
        verloop(scrim, 'bg', 'bg', 5400000)
        for gs4, a, pos in zip(list(scrim.fill._xPr.find(qn('a:gradFill'))
                                    .iter(qn('a:gs'))), (10, 96), ('0', '62000')):
            etree.SubElement(gs4.find(qn('a:srgbClr')), qn('a:alpha')).set(
                'val', str(a * 1000))
            gs4.set('pos', pos)
        foto(s, 'merk.png', 100, 96, 130, 146)
        txt(s, 100, 606, 1400, 260, d['boven'], gr=124, kl='ink', vet=True,
            ra=0.9, omslag=False, naam='hero')
        txt(s, 106, 850, 1400, 60, d['onder'], gr=28, kl='licht', vet=True)
        partners(s, 952)
        txt(s, 1120, 1034, 720, 30, d['voet'], gr=11, kl='gedempt', font=FONT_M,
            uit=PP_ALIGN.RIGHT, omslag=False)

    elif t == 'titel_poster':
        # B — geen beeld: de schaal van de studie is het beeld
        txt(s, 100, 96, 1200, 40, d['boven'] + '   ·   ' + d['onder'], gr=15,
            kl='licht', vet=True, font=FONT_M, caps=True, sp=2.2)
        for j, (g, l) in enumerate(d['tegels']):
            y = 200 + j * 232
            liniaal(s, 100, y, 1720)
            txt(s, 100, y + 26, 460, 220, g, gr=112, kl='licht', vet=True, ra=0.92,
                omslag=False)
            txt(s, 560, y + 96, 1260, 90, l, gr=48, kl='ink', vet=True, ra=1.0,
                omslag=False)
        partners(s, 952)
        txt(s, 1120, 1034, 720, 30, d['voet'], gr=11, kl='gedempt', font=FONT_M,
            uit=PP_ALIGN.RIGHT, omslag=False)

    elif t == 'titel_duo':
        # C — het argument van de studie in één beeld: dezelfde voet, twee keer
        foto(s, 'merk.png', 100, 96, 118, 132)
        txt(s, 260, 108, 900, 60, d['boven'], gr=40, kl='ink', vet=True)
        txt(s, 262, 200, 900, 40, d['onder'], gr=16, kl='gedempt', font=FONT_M)
        for j, (haard, piek, lbl, kl) in enumerate(
                ((beeld.VOOR, 312, 'vandaag', 'oranje'),
                 (beeld.NA, 186, 'met PARADISE', 'licht'))):
            x = 250 + j * 800
            beeld.drukmat(s, x, 268, 440, 560, haard, sleutel='d%d' % j, piek=float(piek))
            txt(s, x - 30, 844, 500, 130, '%d' % piek, gr=64, kl=kl, vet=True,
                ra=1.0, uit=PP_ALIGN.CENTER)
            txt(s, x - 30, 968, 500, 40, 'kPa  ·  ' + lbl, gr=14, kl='gedempt',
                font=FONT_M, uit=PP_ALIGN.CENTER, caps=True, sp=1.4)
        liniaal(s, 1000, 540, 100, 'gedempt', 3)
        txt(s, 100, 1040, 1720, 30, d['voet'], gr=11, kl='gedempt', font=FONT_M)

    elif t == 'titel_vlak':
        # D — één vlak, één woord. Par-dark als grond: institutioneel en luid.
        vlak = tegel(s, 0, 0, 1920, 1080, 'navy', None, rond=False)
        verloop(vlak, 'navy', 'mid', 2700000)
        txt(s, 100, 320, 1740, 340, d['boven'], gr=166, kl='wit', vet=True,
            ra=0.86, omslag=False, naam='hero')
        liniaal(s, 106, 636, 1000, 'wit', 4)
        txt(s, 106, 672, 1400, 60, d['onder'], gr=30, kl='wit', vet=True)
        for j, (g, l) in enumerate(d['tegels']):
            x = 106 + j * 300
            txt(s, x, 820, 280, 80, g, gr=40, kl='wit', vet=True, ra=1.0)
            txt(s, x + 2, 888, 280, 40, l, gr=12.5, kl='wit', font=FONT_M,
                caps=True, sp=1.4)
        partners(s, 952)
        txt(s, 1120, 1034, 720, 30, d['voet'], gr=11, kl='wit', font=FONT_M,
            uit=PP_ALIGN.RIGHT, omslag=False)

    elif t == 'hero_titel':
        # De drukmat zelf is het openingsbeeld: een voet die van de dia loopt,
        # op volle kleur, met de hete plek als brandpunt. Daarover een verloop
        # dat links dichtloopt zodat de titel leest.
        beeld.drukmat(s, 990, -230, 960, 1520, beeld.VOOR, sleutel='t',
                      piek=396.0)
        scrim = tegel(s, 0, 0, 1920, 1080, 'bg', None, rond=False)
        verloop(scrim, 'bg', 'bg', 0)
        sf = scrim.fill._xPr.find(qn('a:gradFill'))
        stops = list(sf.iter(qn('a:gs')))
        for gs2, a, pos in zip(stops, (100, 0), ('0', '78000')):
            etree.SubElement(gs2.find(qn('a:srgbClr')), qn('a:alpha')).set(
                'val', str(a * 1000))
            gs2.set('pos', pos)
        foto(s, 'merk.png', 100, 120, 150, 168, naam='merk')
        txt(s, 100, 336, 1300, 300, d['boven'], gr=132, kl='ink', vet=True,
            ra=0.9, naam='hero', omslag=False)
        txt(s, 106, 590, 1100, 90, d['onder'], gr=40, kl='licht', vet=True)
        txt(s, 108, 682, 1000, 70, d['staart'], gr=18, kl='gedempt', font=FONT_L)
        liniaal(s, 100, 774, 1000)
        for j, (g, l) in enumerate(d['tegels']):
            x = 100 + j * 330
            txt(s, x, 798, 300, 90, g, gr=44, kl='licht', vet=True, ra=1.0)
            txt(s, x + 2, 878, 300, 40, l, gr=12.5, kl='gedempt', font=FONT_M,
                caps=True, sp=1.4)
        # tweede verloop onderaan, zodat de partnerbalk op een rustige grond staat
        voetscrim = tegel(s, 0, 890, 1920, 190, 'bg', None, rond=False)
        verloop(voetscrim, 'bg', 'bg', 5400000)
        sv = voetscrim.fill._xPr.find(qn('a:gradFill'))
        for gs3, a in zip(list(sv.iter(qn('a:gs'))), (0, 94)):
            etree.SubElement(gs3.find(qn('a:srgbClr')), qn('a:alpha')).set(
                'val', str(a * 1000))
        partners(s, 952)
        txt(s, 1120, 1034, 720, 30, d['voet'], gr=11, kl='gedempt', font=FONT_M,
            uit=PP_ALIGN.RIGHT, omslag=False)

    elif t == 'knal':
        # Kleuromkering over de volle dia: één zin, verder niets. Deze dia's
        # breken het ritme en zijn het enige moment waarop de zaal niet leest
        # maar luistert.
        vlak = tegel(s, 0, 0, 1920, 1080, d.get('kleur', 'oranje'), None, rond=False)
        verloop(vlak, d.get('kleur', 'oranje'), d.get('kleur2', 'oranje'), 2700000)
        # elke harde regelovergang is een eigen alinea; die tellen apart mee
        for gk in (96, 84, 72, 62, 54):
            regels = sum(max(1, -(-len(r) // max(1, int((1620 / 2.0) / (gk * 0.5)))))
                         for r in d['kop'].split('\n'))
            if regels * gk * 1.06 * 2.4 <= 620:
                break
        txt(s, 150, 240, 1620, 660, d['kop'], gr=gk, kl='bg', vet=True, ra=1.06,
            naam='!!knal')
        if d.get('onder'):
            txt(s, 154, 946, 1500, 44, d['onder'], gr=17, kl='bg', font=FONT_M)

    elif t == 'sectie':
        # het merk loopt reusachtig van de dia af, en het nummer loopt van de
        # bovenrand: schaal als beeldmiddel in plaats van een nette kop
        foto(s, 'merk.png', 1040, -220, 1240, 1520, alpha=14, naam='merkgroot')
        txt(s, 90, -70, 900, 460, d['nr'], gr=260, vet=True, ra=0.82, naam='hero',
            vloei=('licht', 'oranje'), omslag=False)
        txt(s, 100, 560, 1400, 160, d['titel'], gr=76, kl='ink', vet=True,
            naam='sectietitel')
        txt(s, 104, 760, 1300, 80, d['regel'], gr=23, kl='gedempt', font=FONT_L)
        paginering(s, i)

    elif t == 'keuze':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        for j, (nr, naam, sub) in enumerate(d['tegels']):
            x, w = kol(j * 3, 3)
            tegel(s, x, 340, w, 480, 'tegel', 'rand')
            txt(s, x + 30, 380, w - 60, 66, nr, gr=28, kl='licht', vet=True, font=FONT_M)
            txt(s, x + 30, 466, w - 60, 120, naam, gr=23, kl='ink', vet=True, ra=1.15)
            txt(s, x + 30, 606, w - 60, 180, sub, gr=15, kl='gedempt', ra=1.35)
        paginering(s, i)

    elif t == 'hero_cijfer':
        acc = d.get('accent', 'licht')
        gloed(s, -160, 40, 820, acc, 6)
        # het cijfer loopt bewust buiten het kader: schaal als beeldmiddel
        txt(s, 60, 120, 1200, 520, d['cijfer'], gr=250, kl=acc, vet=True, ra=0.82,
            naam='hero')
        txt(s, 700, 208, 240, 160, d['suffix'], gr=76, kl=acc, vet=True, font=FONT_L)
        txt(s, 104, 626, 860, 200, d['kop'], gr=32, kl='ink', vet=True, ra=1.2)
        if d.get('meter'):
            # meter met streefwaarde: laat zien hoe ver het van de norm ligt
            m = d['meter']
            txt(s, 1000, 300, 800, 36, m['label'], gr=13, kl='gedempt', font=FONT_M,
                caps=True, sp=1.4)
            yy = 350
            for lbl, waarde in m['staven']:
                txt(s, 1000, yy, 260, 36, lbl, gr=15, kl='gedempt')
                staaf(s, 1000, yy + 42, 700, 40, waarde / 100.0, acc)
                txt(s, 1720, yy + 38, 110, 44, '%d%%' % waarde, gr=22, kl='ink', vet=True)
                yy += 118
            merkstreep(s, 1000 + 700 * 0.80, 344, 348, 'oranje')
            txt(s, 1000 + 700 * 0.80 - 100, 704, 200, 36, 'norm 80%', gr=13,
                kl='oranje', vet=True, font=FONT_M, uit=PP_ALIGN.CENTER)
        else:
            for j, (g, l) in enumerate(d['tegels']):
                x, w = kol(8, 4)
                y = 180 + j * 216
                tegel(s, x, y, w, 190, 'tegel', 'rand')
                txt(s, x + 32, y + 28, w - 60, 96, g, gr=42, kl='ink', vet=True, ra=1.0)
                txt(s, x + 34, y + 124, w - 60, 40, l, gr=13.5, kl='gedempt',
                    font=FONT_M, caps=True, sp=1.3)
        txt(s, 100, 1006, 1400, 30, d['voet'], gr=11.5, kl='gedempt', font=FONT_M)
        paginering(s, i)

    elif t == 'splitsing':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        # proportionele staven: lengte draagt de vergelijking, niet het cijfer
        y = 288
        for kant in (d['links'], d['rechts']):
            liniaal(s, 100, y, 1720)
            txt(s, 100, y + 26, 600, 50, kant['titel'], gr=22, kl='ink', vet=True)
            txt(s, 102, y + 82, 600, 34, kant['sub'], gr=13, kl='gedempt', font=FONT_M)
            for lbl, waarde, kl in ((kant['ref_lbl'], kant['ref'], 'gedempt'),
                                    (kant['int_lbl'], float(kant['groot'].replace(',', '.')),
                                     kant['kleur'])):
                yy = y + (140 if kl == 'gedempt' else 206)
                txt(s, 100, yy + 8, 320, 34, lbl, gr=14, kl='gedempt', font=FONT_M)
                staaf(s, 440, yy, 1080, 46, waarde / 50.0, kl)
                txt(s, 1560, yy - 4, 220, 54,
                    ('%.1f' % waarde).replace('.', ',') + '%', gr=27, kl=kl, vet=True)
            txt(s, 100, y + 274, 900, 44, kant['slot'], gr=18, kl=kant['kleur'], vet=True)
            y += 356
        txt(s, 900, 934, 920, 60, d['punchline'], gr=24, kl='ink', vet=True,
            uit=PP_ALIGN.RIGHT)
        txt(s, 100, 1006, 1400, 30, d['voet'], gr=11.5, kl='gedempt', font=FONT_M)
        paginering(s, i)

    elif t == 'formule':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        for j, (naam, sub, kl, tag) in enumerate(d['delen']):
            x, w = kol(j * 4, 4)
            tegel(s, x, 400, w - 40, 310, 'tegel', 'rand')
            txt(s, x + 30, 436, w - 100, 66, naam, gr=25, kl='ink', vet=True)
            txt(s, x + 30, 504, w - 100, 56, sub, gr=15, kl='gedempt')
            tegel(s, x + 30, 598, 160, 44, kl if kl != 'gedempt' else 'glas', None)
            txt(s, x + 30, 610, 160, 32, tag, gr=11,
                kl='bg' if kl != 'gedempt' else 'gedempt', vet=True, font=FONT_M,
                uit=PP_ALIGN.CENTER, caps=True, sp=1.2)
            if j < 2:
                txt(s, x + w - 48, 512, 64, 90, '×', gr=34, kl='gedempt',
                    uit=PP_ALIGN.CENTER, font=FONT_L)
        txt(s, 100, 770, 1720, 90, d['slot'], gr=20, kl='gedempt', ra=1.4)
        paginering(s, i)

    elif t == 'duo':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        for j, c in enumerate((d['een'], d['twee'])):
            x, w = kol(j * 6, 6)
            tegel(s, x, 320, w, 570, 'tegel', 'rand')
            txt(s, x + 40, 354, 200, 60, c['nr'], gr=26, kl='licht', vet=True, font=FONT_M)
            txt(s, x + 40, 418, w - 80, 110, c['naam'], gr=24, kl='ink', vet=True, ra=1.15)
            tegel(s, x + 40, 548, w - 80, 84, 'glas', None)
            txt(s, x + 40, 572, w - 80, 48, c['kern'], gr=18, kl='licht', vet=True,
                uit=PP_ALIGN.CENTER, font=FONT_M)
            txt(s, x + 40, 668, w - 80, 200,
                [(r, {'voor': 13}) for r in c['regels']], gr=15, kl='gedempt', ra=1.35)
        paginering(s, i)

    elif t == 'statraster':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        # 144 deelnemers als puntenmatrix, 72 per arm
        kleuren = ['licht'] * 72 + ['oranje'] * 72
        matrix(s, 100, 300, 144, 24, 26, 12, kleuren)
        txt(s, 100, 520, 400, 40, '72 PARADISE', gr=13, kl='licht', vet=True,
            font=FONT_M, sp=1.4)
        txt(s, 380, 520, 400, 40, '72 GEBRUIKELIJKE ZORG', gr=13, kl='oranje',
            vet=True, font=FONT_M, sp=1.4)
        txt(s, 1080, 286, 740, 250,
            'Elke stip is een patiënt die u zelf includeert.\n24 per kliniek.',
            gr=30, kl='ink', vet=True, ra=1.3)
        for j, (k2, v) in enumerate(d['klein']):
            x, w = kol((j % 2) * 6, 6)
            y = 588 + (j // 2) * 126
            tegel(s, x, y, w, 110, 'tegel', 'rand')
            txt(s, x + 30, y + 18, w - 60, 34, k2, gr=12, kl='licht', vet=True,
                font=FONT_M, caps=True, sp=1.4)
            txt(s, x + 30, y + 56, w - 60, 46, v, gr=15.5, kl='ink')
        # de zes centra bij naam: dat maakt van zes losse klinieken één netwerk
        liniaal(s, 100, 858, 1720)
        txt(s, 100, 878, 900, 34, d['centra_label'], gr=12, kl='gedempt', vet=True,
            font=FONT_M, caps=True, sp=1.4)
        for j, naam2 in enumerate(d['centra']):
            cx = 100 + j * 288
            tegel(s, cx, 926, 276, 52, 'tegel2', 'rand')
            txt(s, cx, 941, 276, 34, naam2, gr=13, kl='ink', vet=True,
                uit=PP_ALIGN.CENTER, omslag=False)
        paginering(s, i)

    elif t == 'steun':
        # asymmetrisch: lijst links, één groot geruststellend cijfer rechts
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        for j, (k2, v) in enumerate(d['rijen']):
            y = 314 + j * 152
            liniaal(s, 100, y, 980)
            txt(s, 100, y + 26, 980, 50, k2, gr=21, kl='ink', vet=True)
            txt(s, 102, y + 76, 960, 60, v, gr=15.5, kl='gedempt', ra=1.3)
        tegel(s, 1180, 300, 640, 622, 'tegel2', 'oranje')
        gloed(s, 1260, 322, 480, 'oranje', 5)
        txt(s, 1224, 340, 560, 300, d['cijfer'], gr=150, kl='oranje', vet=True, ra=0.9,
            uit=PP_ALIGN.CENTER, naam='hero')
        txt(s, 1224, 640, 560, 62, d['cijfer_label'], gr=30, kl='ink', vet=True,
            uit=PP_ALIGN.CENTER)
        txt(s, 1230, 722, 540, 160, d['cijfer_regel'], gr=16, kl='gedempt', ra=1.4,
            uit=PP_ALIGN.CENTER)
        paginering(s, i)

    elif t == 'dubbelkolom':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        for j, kant in enumerate((d['links'], d['rechts'])):
            x, w = kol(j * 6, 6)
            tegel(s, x, 314, w, 606, 'tegel', 'rand')
            tegel(s, x + 40, 348, 120, 46, kant['kleur'], None)
            txt(s, x + 40, 360, 120, 34, kant['titel'], gr=13.5, kl='bg', vet=True,
                uit=PP_ALIGN.CENTER, font=FONT_M, caps=True, sp=1.3)
            txt(s, x + 40, 432, w - 80, 470,
                [(it, {'voor': 12}) for it in kant['items']], gr=14, kl='gedempt',
                ra=1.3)
        paginering(s, i)

    elif t in ('traject', 'zoom'):
        # Doorlopend canvas: elke dia toont hetzelfde traject onder een andere
        # camera-instelling. Morph koppelt de vormen via hun !!-naam en maakt
        # er een vloeiende in- of uitzoom van in plaats van een dia-wissel.
        kicker(s, d['kicker'], y=56)
        kop(s, d['kop'], y=96, naam='!!sectietitel')
        band = tegel(s, 0, 190, 1920, 520, 'tegel', None, rond=False)
        band.name = '!!band'
        if t == 'traject':
            beeld.canvas(s, PUNTEN, txt, liniaal, 1.0, 960, 454, 960, 460)
            beeld.legende(s, 104, 632, SOORT, txt)
            txt(s, 100, 742, 320, 230, d['groot'], gr=96, kl='licht', vet=True,
                ra=1.0, omslag=False, naam='!!camera')
            txt(s, 400, 758, 1400, 66, d['punt'], gr=30, kl='ink', vet=True)
            txt(s, 402, 832, 1400, 150, d['slot'], gr=18, kl='gedempt', ra=1.4)
        else:
            beeld.canvas(s, PUNTEN, txt, liniaal, d['schaal'], d['mid'], 454,
                         d.get('ox', 560), 500)
            p = d['paneel']
            liniaal(s, 100, 748, 1720)
            txt(s, 100, 772, 420, 60, p['titel'], gr=24, kl='oranje', vet=True,
                naam='!!camera')
            for j, r in enumerate(p['regels']):
                txt(s, 560 + j * 660, 770, 600, 200, r, gr=16.5, kl='gedempt', ra=1.35)
        paginering(s, i)

    elif t == 'doel':
        # de leerdoelen van de dag; ze komen op het eind terug als toets
        kicker(s, d['kicker']); kop(s, d['kop'], naam='!!sectietitel')
        for j, (naam2, uitleg) in enumerate(d['doelen']):
            x = 100 + (j % 2) * 900
            y = 288 + (j // 2) * 152
            liniaal(s, x, y, 820)
            bol(s, x, y + 22, 38, str(j + 1), 'licht')
            txt(s, x + 58, y + 18, 760, 48, naam2, gr=21, kl='ink', vet=True)
            txt(s, x + 58, y + 70, 760, 80, uitleg, gr=14.5, kl='gedempt', ra=1.3)
        tegel(s, 100, 764, 1720, 116, 'tegel', 'licht')
        txt(s, 150, 792, 1620, 70, d['slot'], gr=19, kl='ink', ra=1.4)
        paginering(s, i)

    elif t == 'meetreeks':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='!!sectietitel')
        for j, (nr, naam2, uitleg) in enumerate(d['condities']):
            y = 292 + j * 176
            tegel(s, 100, y, 1720, 152, 'tegel', 'rand')
            txt(s, 148, y + 34, 120, 70, nr, gr=34, kl='licht', vet=True, font=FONT_M)
            txt(s, 300, y + 30, 620, 54, naam2, gr=24, kl='ink', vet=True)
            txt(s, 960, y + 30, 812, 100, uitleg, gr=15.5, kl='gedempt', ra=1.32)
        txt(s, 100, 848, 1720, 90, d['slot'], gr=19, kl='oranje', vet=True, ra=1.4)
        txt(s, 100, 1006, 1500, 30, d['voet'], gr=11.5, kl='gedempt', font=FONT_M)
        paginering(s, i)

    elif t == 'bezoek' and not d.get('kaart'):
        # compacte versie: alles op één dia, voor de decks zonder stapreeks
        kicker(s, d['kicker']); kop(s, d['kop'], naam='!!sectietitel')
        tegel(s, 100, 236, 560, 44, 'glas', None)
        txt(s, 100, 247, 560, 32, d['wanneer'], gr=13, kl='licht', vet=True,
            uit=PP_ALIGN.CENTER, font=FONT_M)
        y = 322
        for j, h in enumerate(d['handelingen']):
            h = h[0] if isinstance(h, tuple) else h
            bol(s, 100, y, 30, str(j + 1), 'licht', klein=True)
            hh = hoogte(h, 1020, 13, 1.28)
            txt(s, 146, y - 2, 1020, hh + 10, h, gr=13, kl='ink', ra=1.28)
            y += hh + 22
        txt(s, 1220, 322, 600, 34, 'Wat u invult', gr=12, kl='gedempt', vet=True,
            font=FONT_M, caps=True, sp=1.5)
        nrs = '   '.join(nr for nr, _ in d['documenten'])
        txt(s, 1220, 366, 600, 120, nrs, gr=19, kl='licht', vet=True, font=FONT_M,
            ra=1.4)
        # PowerPoint zet regels iets ruimer dan de schatting; vandaar de marge
        for gl in (14, 13, 12, 11):
            hl = hoogte(d['letop'], 520, gl, 1.3) * 1.22
            if hl <= 344:
                break
        tegel(s, 1220, 520, 600, hl + 96, 'tegel', 'oranje')
        txt(s, 1258, 542, 520, 34, 'Let op', gr=12, kl='oranje', vet=True,
            font=FONT_M, caps=True, sp=1.5)
        txt(s, 1258, 580, 520, hl + 20, d['letop'], gr=gl, kl='ink', ra=1.3)
        txt(s, 100, 1006, 1500, 30, d['wie'], gr=11.5, kl='gedempt', font=FONT_M)
        paginering(s, i)

    elif t == 'bezoek':
        # titelkaart van het bezoek; de handelingen komen erna, één per dia
        gloed(s, 1180, 180, 760, 'licht', 5)
        kicker(s, d['kicker'])
        ht = hoogte(d['kop'], 1100, 56, 1.1)
        txt(s, 100, 300, 1100, ht + 40, d['kop'], gr=56, kl='ink', vet=True, ra=1.1,
            naam='!!sectietitel')
        yc = 300 + ht + 84
        tegel(s, 100, yc, 560, 48, 'glas', None)
        txt(s, 100, yc + 12, 560, 34, d['wanneer'], gr=14, kl='licht', vet=True,
            uit=PP_ALIGN.CENTER, font=FONT_M)
        for j, (g, l) in enumerate(((str(len(d['handelingen'])), 'handelingen'),
                                    (str(len(d['documenten'])), 'formulieren'))):
            x = 1260 + j * 300
            txt(s, x, 560, 260, 130, g, gr=64, kl='licht', vet=True, ra=1.0)
            txt(s, x + 2, 690, 260, 40, l, gr=13.5, kl='gedempt', font=FONT_M,
                caps=True, sp=1.4)
        liniaal(s, 100, 900, 1720)
        txt(s, 100, 924, 1720, 40, d['wie'], gr=15, kl='gedempt')
        paginering(s, i)

    elif t == 'stap':
        # keynotetempo: één handeling per dia, groot, met een voortgangsrail
        kicker(s, d['visite'])
        txt(s, 100, 236, 340, 300, '%02d' % d['nr'], gr=150, kl='licht', vet=True,
            ra=0.9, omslag=False, uit=PP_ALIGN.RIGHT, naam='!!stapnr')
        if d.get('bijschrift'):
            # tweekolommig: handeling links, fotokader rechts om later te vullen
            tw, ruimte = 620, 480
            foto(s, 'plaatshouder.png', 1120, 250, 700, 440, naam='fotokader')
            txt(s, 1120, 706, 700, 60, d['bijschrift'], gr=13.5, kl='gedempt',
                font=FONT_M)
        else:
            tw = 1320
            ruimte = 330 if d['documenten'] else 560
        for gs in (40, 36, 32, 28, 24):
            if hoogte(d['tekst'], tw, gs, 1.25) <= ruimte:
                break
        txt(s, 500, 250, tw, ruimte + 20, d['tekst'], gr=gs, kl='ink', vet=True,
            ra=1.25, naam='!!staptekst')
        if d['documenten']:
            txt(s, 470, 640, 1350, 34, 'Wat u hierbij invult', gr=12.5, kl='gedempt',
                vet=True, font=FONT_M, caps=True, sp=1.5)
            for j, (nr, naam2) in enumerate(d['documenten']):
                x = 470 + (j % 5) * 272
                y = 686 + (j // 5) * 54
                tegel(s, x, y, 262, 44, 'tegel2', 'rand')
                txt(s, x + 12, y + 12, 56, 30, nr, gr=12, kl='licht', vet=True,
                    font=FONT_M)
                txt(s, x + 66, y + 14, 190, 28, naam2, gr=10.5, kl='ink', omslag=False)
        if d['letop']:
            liniaal(s, 100, 858, 1720)
            txt(s, 100, 880, 200, 34, 'Let op', gr=12.5, kl='oranje', vet=True,
                font=FONT_M, caps=True, sp=1.5)
            txt(s, 320, 876, 1500, 90, d['letop'], gr=15.5, kl='ink', ra=1.3)
        # rail onderaan: waar zit u in de reeks
        bw = 1720.0 / d['totaal']
        for k in range(d['totaal']):
            aan = k == d['nr'] - 1
            r = liniaal(s, 100 + k * bw + 4, 1006, bw - 8,
                        'licht' if aan else 'rand', 6 if aan else 3)
            r.name = '!!rail%d' % k
        paginering(s, i)

    elif t == 'sop':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='!!sectietitel')
        txt(s, 100, 300, 820, 34, 'Instellingen', gr=12.5, kl='gedempt', vet=True,
            font=FONT_M, caps=True, sp=1.5)
        y = 348
        for sleutel, waarde in d['instellingen']:
            liniaal(s, 100, y, 820)
            hw = hoogte(waarde, 480, 14, 1.25)
            txt(s, 100, y + 16, 330, 34, sleutel, gr=13.5, kl='gedempt', omslag=False)
            txt(s, 440, y + 14, 480, hw + 10, waarde, gr=14, kl='ink', vet=True, ra=1.25)
            y += max(64, hw + 34)
        einde_links = y
        txt(s, 1000, 300, 820, 34, 'Stappen', gr=12.5, kl='gedempt', vet=True,
            font=FONT_M, caps=True, sp=1.5)
        y = 348
        for j, stap in enumerate(d['stappen']):
            bol(s, 1000, y + 2, 32, str(j + 1), 'licht', klein=True)
            hs = hoogte(stap, 772, 13.5)
            txt(s, 1048, y, 772, hs + 10, stap, gr=13.5, kl='ink', ra=1.3)
            y += hs + 30
        yv = max(einde_links, y) + 24
        hv = hoogte(d['valkuil'], 1620, 15, 1.35)
        tegel(s, 100, yv, 1720, hv + 96, 'tegel', 'oranje')
        txt(s, 148, yv + 22, 300, 34, 'Valkuil', gr=12.5, kl='oranje', vet=True,
            font=FONT_M, caps=True, sp=1.5)
        txt(s, 148, yv + 62, 1620, hv + 20, d['valkuil'], gr=15, kl='ink', ra=1.35)
        paginering(s, i)

    elif t == 'melden':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='!!sectietitel')
        for j, (wat, hoe, nrs) in enumerate(d['rijen']):
            y = 296 + j * 132
            liniaal(s, 100, y, 1720)
            txt(s, 100, y + 24, 380, 50, wat, gr=20, kl='ink', vet=True)
            txt(s, 520, y + 24, 980, 90, hoe, gr=15, kl='gedempt', ra=1.32)
            txt(s, 1560, y + 26, 260, 40, nrs, gr=13, kl='licht', vet=True,
                font=FONT_M, uit=PP_ALIGN.RIGHT)
        tegel(s, 100, 846, 1720, 116, 'tegel', 'licht')
        txt(s, 150, 874, 1620, 70, d['slot'], gr=19, kl='ink', ra=1.4)
        paginering(s, i)

    elif t == 'check':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='!!sectietitel')
        for j, it in enumerate(d['items']):
            x = 100 + (j % 2) * 900
            y = 292 + (j // 2) * 122
            vk = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y + 8),
                                    px(38), px(38))
            vk.fill.solid(); vk.fill.fore_color.rgb = rgb('bg')
            vk.line.color.rgb = rgb('licht'); vk.line.width = Pt(1.5)
            vk.shadow.inherit = False; vk.adjustments[0] = 0.2
            txt(s, x + 58, y + 4, 780, 100, it, gr=15.5, kl='ink', ra=1.3)
        tegel(s, 100, 794, 1720, 150, 'tegel', 'oranje')
        txt(s, 150, 822, 1620, 110, d['slot'], gr=18, kl='ink', ra=1.4)
        paginering(s, i)

    elif t == 'register':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='!!sectietitel')
        tegel(s, 100, 296, 760, 400, 'tegel2', 'rand')
        txt(s, 144, 336, 672, 330, d['body'], gr=17, kl='ink', ra=1.45)
        for j, (wanneer, wat) in enumerate(d['regels']):
            y = 296 + j * 138
            liniaal(s, 920, y, 900)
            txt(s, 920, y + 22, 900, 44, wanneer, gr=20, kl='licht', vet=True)
            txt(s, 920, y + 66, 900, 70, wat, gr=15, kl='gedempt', ra=1.3)
        txt(s, 100, 744, 1720, 70, d['slot'], gr=19, kl='ink', vet=True, ra=1.4)
        paginering(s, i)

    elif t == 'aflopend':
        # Aflopend beeld tot in de vier randen, met een verloop eroverheen zodat
        # de tekst leesbaar blijft. Eén regel, verder niets.
        foto(s, d['foto'], 0, 0, 1920, 1080, vullend=True, naam='!!beeld')
        # verloop van doorzichtig bovenaan naar bijna dicht onderaan, zodat de
        # tekst leesbaar wordt zonder het beeld af te dekken
        scrim = tegel(s, 0, 0, 1920, 1080, 'bg', None, rond=False)
        verloop(scrim, 'bg', 'bg', 5400000)
        sf = scrim.fill._xPr.find(qn('a:gradFill'))
        for gs, a in zip(sf.iter(qn('a:gs')), (6, 94)):
            etree.SubElement(gs.find(qn('a:srgbClr')), qn('a:alpha')).set(
                'val', str(a * 1000))
        kicker(s, d['kicker'], y=706)
        txt(s, 100, 754, 1500, 240, d['kop'], gr=44, kl='ink', vet=True, ra=1.15,
            naam='!!sectietitel')
        paginering(s, i)

    elif t == 'drukzoom':
        # Zelfde cellen, zelfde !!-namen, alleen een ander kader: de camera
        # duikt de voorvoet in. Morph maakt daar een inzoom van.
        # Zuivere schaalsprong van 2,6× ten opzichte van dia 9: de cellen houden
        # hun onderlinge verhouding, dus de camera zoomt in plaats van te rekken.
        MX, MY, MW, MH = -85, 240, 1170, 1794
        beeld.drukmat(s, MX, MY, MW, MH,
                      beeld.NA if d['fase'] == 'na' else beeld.VOOR, sleutel='p',
                      piek=float(d['waarde']))
        kicker(s, d['kicker'])
        kop(s, d['kop'], naam='!!sectietitel')
        for j, (naam2, uitleg, fx, ft) in enumerate(d['regios']):
            # genummerde speld op de voet, met hetzelfde nummer in de lijst
            cx, cy = MX + fx * MW, MY + MH * (1 - ft)
            bol(s, cx - 22, cy - 22, 44, str(j + 1), 'wit', vul='bg')
            y = 300 + j * 116
            liniaal(s, 1000, y, 820)
            bol(s, 1000, y + 26, 40, str(j + 1), 'licht')
            txt(s, 1064, y + 22, 760, 44, naam2, gr=20, kl='ink', vet=True)
            txt(s, 1066, y + 64, 760, 40, uitleg, gr=14.5, kl='gedempt')
        txt(s, 100, 1006, 1400, 30, d['voet'], gr=11.5, kl='gedempt', font=FONT_M)
        paginering(s, i)

    elif t == 'drukmeting':
        # De sensormatrix is geen illustratie maar de meting zelf: dezelfde
        # cellen op beide dia's, alleen de kleur verandert. Morph laat de
        # drukpiek daardoor letterlijk afkoelen.
        kicker(s, d['kicker'])
        kop(s, d['kop'], naam='!!sectietitel')
        beeld.drukmat(s, 70, 216, 450, 690,
                      beeld.NA if d['fase'] == 'na' else beeld.VOOR, sleutel='p',
                      piek=float(d['waarde']))
        beeld.schaalbalk(s, 145, 930, 300, 16, txt)
        acc = 'licht' if d['fase'] == 'na' else 'oranje'
        txt(s, 560, 250, 700, 300, d['waarde'], gr=150, kl=acc, vet=True, ra=0.9,
            omslag=False, naam='!!piek')
        txt(s, 1080, 300, 200, 80, 'kPa', gr=44, kl=acc, vet=True, font=FONT_L,
            naam='!!eenheid')
        txt(s, 564, 528, 1200, 60, d['plek'], gr=22, kl='ink', vet=True,
            naam='!!plek')
        tegel(s, 560, 618, 1260, 96, 'glas', None, naam='!!norm')
        txt(s, 560, 646, 1260, 46, d['norm'], gr=22, kl='oranje', vet=True,
            uit=PP_ALIGN.CENTER, font=FONT_M, naam='!!normtekst')
        txt(s, 560, 760, 1260, 200,
            [(r, {'voor': 12}) for r in d['regels']], gr=17, kl='gedempt', ra=1.35)
        txt(s, 100, 1006, 1400, 30, d['voet'], gr=11.5, kl='gedempt', font=FONT_M)
        paginering(s, i)

    elif t == 'vijfluik':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        # de stappen staan op één lijn: het ís een reeks, dus laat dat zien
        liniaal(s, 100, 388, 1720, 'rand', 2)
        for j, (nr, wanneer, wat) in enumerate(d['stappen']):
            x = 100 + j * 348
            o = s.shapes.add_shape(MSO_SHAPE.OVAL, px(x + 18), px(370), px(36), px(36))
            o.fill.solid(); o.fill.fore_color.rgb = rgb('licht')
            o.line.color.rgb = rgb('bg'); o.line.width = Pt(3); o.shadow.inherit = False
            tegel(s, x, 440, 324, 400, 'tegel', 'rand')
            txt(s, x + 28, 474, 260, 60, nr, gr=27, kl='licht', vet=True, font=FONT_M)
            txt(s, x + 28, 546, 268, 66, wanneer, gr=20, kl='ink', vet=True)
            txt(s, x + 28, 628, 268, 190, wat, gr=15, kl='gedempt', ra=1.38)
        paginering(s, i)

    elif t == 'drieluik':
        # echte bento: één groot vlak links, twee gestapelde rechts
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        g, rest = d['kaarten'][0], d['kaarten'][1:]
        gx, gw = kol(0, 6)
        tegel(s, gx, 310, gw, 630, 'tegel', 'rand')
        foto(s, g['foto'], gx + 40, 344, gw - 80, 262)
        txt(s, gx + 44, 626, 380, 70, g['naam'], gr=32, kl='ink', vet=True)
        txt(s, gx + 46, 700, 380, 36, g['sub'], gr=12.5, kl='licht', vet=True,
            font=FONT_M, caps=True, sp=1.6)
        txt(s, gx + 44, 734, 380, 60, g['groot'] + '  ' + g['onder'], gr=24,
            kl='licht', vet=True, ra=1.0)
        txt(s, gx + 44, 800, 380, 40, g['wie'], gr=13, kl='oranje', vet=True,
            font=FONT_M, caps=True, sp=1.3)
        txt(s, gx + 448, 626, gw - 492, 280,
            [(r, {'voor': 10}) for r in g['regels']], gr=14.5, kl='gedempt', ra=1.3)
        for j, c in enumerate(rest):
            x, w = kol(6, 6)
            y = 310 + j * 328
            tegel(s, x, y, w, 302, 'tegel', 'rand')
            tegel(s, x + 24, y + 24, 340, 254, c.get('vlak', 'tegel2'), None)
            foto(s, c['foto'], x + 36, y + 36, 316, 230)
            tx = x + 400
            txt(s, tx, y + 30, 420, 54, c['naam'], gr=25, kl='ink', vet=True)
            txt(s, tx + 2, y + 88, 420, 34, c['sub'], gr=12, kl='licht', vet=True,
                font=FONT_M, caps=True, sp=1.5)
            txt(s, tx, y + 120, 420, 56, c['groot'] + '  ' + c['onder'], gr=22,
                kl='licht', vet=True, ra=1.0)
            txt(s, tx, y + 174, w - 440, 90,
                [(r, {'voor': 6}) for r in c['regels']], gr=14, kl='gedempt', ra=1.28)
            txt(s, tx, y + 264, 420, 36, c['wie'], gr=12.5, kl='oranje', vet=True,
                font=FONT_M, caps=True, sp=1.3)
        paginering(s, i)

    elif t == 'vierluik':
        # rijen in plaats van kolommen: leest sneller en breekt het kaartenritme
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        for j, (rol, punten) in enumerate(d['kolommen']):
            y = 304 + j * 164
            liniaal(s, 100, y, 1720)
            txt(s, 100, y + 34, 420, 60, rol, gr=27, kl='ink', vet=True)
            for k, p in enumerate(punten):
                px_ = 560 + k * 428
                punt(s, px_, y + 48, 11, 'licht')
                txt(s, px_ + 26, y + 36, 372, 90, p, gr=15.5, kl='gedempt', ra=1.3)
        liniaal(s, 100, 960, 1720)
        paginering(s, i)

    elif t == 'statement':
        kicker(s, d['kicker'], kl='oranje')
        txt(s, 100, 168, 1080, 400, d['kop'], gr=54, kl='ink', vet=True, ra=1.08, naam='hero')
        tegel(s, 100, 646, 1080, 244, 'tegel', 'rand')
        txt(s, 144, 690, 1000, 190, d['body'], gr=17.5, kl='gedempt', ra=1.45)
        tegel(s, 1240, 320, 580, 554, 'tegel2', 'oranje')
        txt(s, 1288, 366, 490, 38, 'Uitzondering', gr=12.5, kl='oranje', vet=True,
            font=FONT_M, caps=True, sp=1.6)
        txt(s, 1288, 430, 490, 400, d['uitzondering'], gr=18, kl='ink', ra=1.45)
        paginering(s, i)

    elif t == 'eerlijk':
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        txt(s, 100, 264, 1400, 44, d['intro'], gr=16.5, kl='gedempt')
        for j, (k2, v) in enumerate(d['rijen']):
            x, w = kol((j % 2) * 6, 6)
            y = 336 + (j // 2) * 142
            tegel(s, x, y, w, 122, 'tegel', 'rand')
            txt(s, x + 32, y + 22, w - 64, 36, k2, gr=12.5, kl='licht', vet=True,
                font=FONT_M, caps=True, sp=1.4)
            txt(s, x + 32, y + 64, w - 64, 46, v, gr=17, kl='ink')
        tegel(s, 100, 646, 1720, 150, 'tegel2', 'licht')
        txt(s, 150, 692, 1620, 110, d['slot'], gr=19, kl='ink', ra=1.4)
        paginering(s, i)

    elif t == 'pauze':
        kicker(s, d['kicker'], kl='oranje')
        txt(s, 100, 372, 1720, 280, d['kop'], gr=120, kl='ink', vet=True,
            uit=PP_ALIGN.CENTER, ra=1.0, naam='hero')
        txt(s, 100, 700, 1720, 70, d['regel'], gr=23, kl='gedempt',
            uit=PP_ALIGN.CENTER, font=FONT_L)
        paginering(s, i)

    elif t == 'vraagraster':
        # geen kaders: enkel haarlijnen, vraag in wit, antwoord ingesprongen
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        for j, (v, a) in enumerate(d['paren']):
            x, w = kol((j % 2) * 6, 6)
            y = 300 + (j // 2) * 226
            liniaal(s, x, y, w)
            txt(s, x, y + 26, w, 60, v, gr=19, kl='ink', vet=True, ra=1.2)
            txt(s, x, y + 96, 30, 40, '—', gr=15, kl='licht', vet=True)
            txt(s, x + 40, y + 94, w - 40, 100, a, gr=15.5, kl='gedempt', ra=1.35)
        paginering(s, i)

    elif t == 'afspraak':
        # één blad in plaats van vier tegels: dit leest als een overeenkomst
        kicker(s, d['kicker']); kop(s, d['kop'], naam='sectietitel')
        tegel(s, 100, 300, 1720, 622, 'tegel', 'rand')
        for j, (g, m, sub) in enumerate(d['tegels']):
            y = 336 + j * 150
            if j:
                liniaal(s, 168, y - 26, 1584)
            txt(s, 118, y, 300, 120, g, gr=54, kl='licht', vet=True, ra=1.0,
                uit=PP_ALIGN.RIGHT, omslag=False)
            txt(s, 470, y, 900, 60, m, gr=26, kl='ink', vet=True)
            txt(s, 472, y + 66, 1250, 50, sub, gr=16, kl='gedempt', ra=1.0)
        paginering(s, i)

    elif t == 'slot':
        txt(s, 100, 300, 1200, 230, d['kop'], gr=110, kl='ink', vet=True, ra=1.0,
            naam='hero')
        txt(s, 104, 566, 1400, 80, d['regel'], gr=25, kl='licht', font=FONT_L)
        partners(s, 900)
        txt(s, 100, 1022, 1500, 30, d['voet'], gr=11.5, kl='gedempt', font=FONT_M)

    elif t == 'contact':
        # De dia die blijft staan tijdens het napraten: gezicht erbij, zodat
        # mensen weten wie ze moeten aanspreken.
        kicker(s, d['kicker'])
        kop(s, d['kop'], naam='!!sectietitel')
        tegel(s, 100, 300, 500, 620, 'tegel2', 'rand')
        foto(s, 'janou.png', 116, 316, 468, 588, naam='portret')
        x = 680
        for j, p in enumerate(d['personen']):
            y = 300 + j * 310
            liniaal(s, x, y, 1140)
            txt(s, x, y + 26, 1000, 60, p['naam'], gr=30, kl='ink', vet=True)
            txt(s, x + 2, y + 100, 900, 36, p['rol'], gr=13, kl='licht', vet=True,
                font=FONT_M, caps=True, sp=1.5)
            txt(s, x, y + 150, 1140, 100, p['waarvoor'], gr=16.5, kl='gedempt', ra=1.35)
            for k, regel in enumerate(p['bereik']):
                txt(s, x + k * 560, y + 256, 580, 40, regel, gr=16, kl='ink',
                    font=FONT_M, omslag=False)
        liniaal(s, 100, 940, 1720)
        txt(s, 100, 962, 1720, 60, d['consortium'], gr=13, kl='gedempt', ra=1.35)
        paginering(s, i)

    else:
        raise SystemExit('onbekend type: %s' % t)

    merk(s, t)
    notitie(s, d)
    morph(s, d.get('morph', 'morph'))

prs.save(DOEL)
print("%d dia's · %.0f KB · %s" % (TOT, os.path.getsize(DOEL) / 1024, os.path.basename(DOEL)))
